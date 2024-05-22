#!/usr/bin/env python
# coding: utf-8

import os
import subprocess
import re
import pandas as pd
from openpyxl import load_workbook
import openpyxl
from openpyxl.utils.dataframe import dataframe_to_rows

import numpy as np
import statsmodels.formula.api as stats
import matplotlib.pyplot as plt
import os
import time
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.dml import MSO_THEME_COLOR
import psutil

# Initialize data dictionaries
data_dict = {}
data = []
logdir=os.getcwd()

# Create Error Folder if not exists
error_folder_path = os.path.join(logdir, "Error Folder")
if not os.path.exists(error_folder_path):
    os.makedirs(error_folder_path)

def time_value(string):
    return string.split("-")[0]+string.split("-")[1]+string.split("-")[2]

def date_stamp(string):
    return string.split("-")[0]+"/"+string.split("-")[1]+"/"+string.split("-")[2]+"-"

def clock_stamp(string):
    return string.split("-")[0]+":"+string.split("-")[1]+":"+string.split("-")[2]

def extract_reading(sensor_string, pattern):
    # Define the regular expression pattern to match the pressure reading
    match = re.search(pattern, sensor_string)

    if match:
        light_reading = float(match.group(1))
        return light_reading
    else:
        # Return None if no match is found
        return None

def update_dict_in_array(array_of_dicts, new_object):
    for i, old_dict in enumerate(array_of_dicts):
        if old_dict["Unit ID"] == new_object["Unit ID"]:
            array_of_dicts[i] = new_object
            break

for file in os.listdir(logdir+"\\logs\\"):
    # Get unit id
    unit_id = file.split('_')[0]
    # So we can take the newer one
    time_val = int(time_value(file.split('_')[1])+time_value(file.split('_')[2]))
    time_stamp = date_stamp(file.split('_')[1])+clock_stamp(file.split('_')[2])
    file_path = "logs\\" + file
    f = open(file_path, 'r')
    try:
        lines = f.readlines()
        button_pushed = False
        wifi_scan = 0
        version = "None"
        at_read_1_found = False
        scan_record = 0
        light = []
        pressure = []
        accXYZ = []
        for line in lines:
            if re.search('Quick charge test with battery:', line):
                string=line.split(":")
                fast_charge_current = float(string[1].split("\n")[0])
            if re.search('IMEI:', line):
                string=line.split(":")
                IMEI = int(string[1].split("=")[0])
            if re.search('Battery voltage ', line):
                string=line.split(" ")
                battery_voltage = int(string[2])
            if re.search('%CCID: ', line):
                string=line.split(":")
                CCID = int(string[1])
            if re.search('>> Temp Record', line):
                string=line.split(":")
                temp_sensor = float(string[3].split(" ")[1])
            if re.search('>> BLE ping response len: ', line):
                string=line.split(":")
                ble = int(string[1].split(" ")[1].split(",")[0])
            if re.search('Button pushed', line):
                button_pushed = True
            if re.search('Bin version:', line):
                wifi_connection_version = line.split("(")[1].split("-")[0]
                wifi_connection = 'ESP32C3' in line
            if re.search(r'Record (\d+): \+CWLAP',line):
                wifi_scan += 1
            if re.search('blename:', line):
                broadcast_SN = line.split(':')[1].split(',')[0]
                broadcast_SN_pass = 'ESP_' in line
            if re.search('SMM3', line):
                version = line.split('\n')[0]
            if re.search('send commond>at+read 1', line):
                at_read_1_found = True
            if re.search('%IGNSSINFO: ', line):
                satellite = int(line.split(" ")[1].split("\n")[0])
            if re.search('Scan Record: ', line):
                scan_record += 1
            if re.search(">> Sensor Record", line):
                light.append(int(extract_reading(line, r'Light: (\d+) lux')))
                pressure.append(float(extract_reading(line, r'Pressure: (\d+\.\d+) hPa')))
                accXYZ.append([float(extract_reading(line, r'accX: ([\d\.-]+)')), 
                               float(extract_reading(line, r'accY: ([\d\.-]+)')), 
                               float(extract_reading(line, r'accZ: ([\d\.-]+)'))])
            if re.search("Read bsn:", line):
                SN_matches_id = file.split('_')[0] == line.split(':')[1].split('\n')[0]

        # Create a dictionary for current log file
        data_dict = {
            "Unit ID": unit_id,
            "SN matches Unit ID": SN_matches_id,
            "Timestamp": time_stamp,
            "GPIB Current Reading": fast_charge_current,
            "GPIB Current Reading Pass/Fail": 430 <= fast_charge_current and 860 >= fast_charge_current,
            "IMEI": IMEI,
            "Battery voltage": battery_voltage,
            "Battery voltage Pass/Fail": 4600 <= battery_voltage and 6200 >= battery_voltage,
            "CCID": CCID,
            "Version": version,
            "Satellite": satellite,
            "Satellite Pass/Fail": satellite>=1,
            "Min Light": min(light),
            "Max Light": max(light),
            "First Pressure": pressure[0],
            "Last Pressure": pressure[len(pressure)-1],
            "First AccX": accXYZ[0][0],
            "First AccY": accXYZ[0][1],
            "First AccZ": accXYZ[0][2],
            "Last AccX": accXYZ[len(accXYZ)-1][0],
            "Last AccY": accXYZ[len(accXYZ)-1][1],
            "Last AccZ": accXYZ[len(accXYZ)-1][2],
            "Temp Sensor": temp_sensor,
            "Temp Sensor Pass/Fail": 20 <= temp_sensor and 30 >= temp_sensor,
            "BLE": ble,
            "BLE Pass/Fail": 20 == ble,
            "Scan Record": scan_record,
            "Button Pushed": button_pushed,
            "WiFi Connection Version": wifi_connection_version,
            "WiFi Connection": wifi_connection,
            "WiFi Scan": wifi_scan,
            "Broadcast SN": broadcast_SN,
            "Broadcast SN Pass/Fail": broadcast_SN_pass,
            "Time Value": time_val
        }

        # Check for duplicate IDs
        contained = False
        replace = False
        for unit in data:
            if unit["Unit ID"] == unit_id:
                contained = True
                if unit["Time Value"] < time_val:
                    replace = True
        if not contained:
            data.append(data_dict)
        else:
            if replace:
                update_dict_in_array(data, data_dict)

    except Exception as e:
        f.close()
        # Move this file to directory "Error Folder"
        print(f"Error occurred on file: {file}. {str(e)}")
        error_file_path = os.path.join(error_folder_path, file)
        os.rename(file_path, error_file_path)
print(len(data))

# Create a dataframe
df = pd.DataFrame(data)
df = df[['Unit ID','SN matches Unit ID', 'Timestamp','GPIB Current Reading', 'GPIB Current Reading Pass/Fail', 'IMEI', 
         'Battery voltage', 'Battery voltage Pass/Fail', 'CCID', 'Version', 'Satellite', 'Satellite Pass/Fail',
         'Min Light', 'Max Light', 'First Pressure', 'Last Pressure', 'First AccX', 'First AccY', 'First AccZ', 
         'Last AccX', 'Last AccY', 'Last AccZ', 'Temp Sensor', 'Temp Sensor Pass/Fail', 'BLE', 'BLE Pass/Fail', 'Scan Record', 
         'Button Pushed', 'WiFi Connection Version', 'WiFi Connection', 'WiFi Scan', 'Broadcast SN', 'Broadcast SN Pass/Fail', 'Time Value']]
# Save dataframe as excel
if os.path.isfile(logdir+'\\ASSY-MMI-Summary.xlsx') == False:
    df.to_excel(logdir+"\\ASSY-MMI-Summary.xlsx", index=False)
else:
    workbook = openpyxl.load_workbook(logdir+'\\ASSY-MMI-Summary.xlsx')  # load workbook if already exists
    sheet = workbook['Sheet1']  # declare the active sheet 
    # append the dataframe results to the current excel file
    for row in dataframe_to_rows(df, header = False, index = False):
        sheet.append(row)
    workbook.save(logdir+'\\ASSY-MMI-Summary.xlsx')  # save workbook
    workbook.close()  # close workbook

stats = ['GPIB Current Reading','Battery voltage','Satellite','Min Light','Max Light','First Pressure',
         'First AccX','First AccY','First AccZ','Temp Sensor','WiFi Scan']
prs = Presentation()

# Before trying to go to a specific slide, ensure PowerPoint is open and there is an active presentation
def wait_for_powerpoint():
    while "POWERPNT.EXE" not in (p.name() for p in psutil.process_iter()):
        time.sleep(1)

# Function to calculate pass/fail statistics for a given metric
def calculate_pass_fail(data, metric):
    if metric == "GPIB Current Reading":
        pass_count = sum(1 for entry in data if entry[metric] >= 430 and entry[metric] <= 860)
        fail_count = sum(1 for entry in data if entry[metric] < 430 or entry[metric] > 860)
    elif metric == "Battery voltage":
        pass_count = sum(1 for entry in data if entry[metric] >= 4600 and entry[metric] <= 6200)
        fail_count = sum(1 for entry in data if entry[metric] < 4600 or entry[metric] > 6200)
    elif metric == "Satellite":
        pass_count = sum(1 for entry in data if entry[metric] >= 1)
        fail_count = sum(1 for entry in data if entry[metric] < 1)
    elif metric == "Min Light":
        pass_count = sum(1 for entry in data if entry[metric] <= 5)
        fail_count = sum(1 for entry in data if entry[metric] > 5)
    elif metric == "Max Light":
        pass_count = sum(1 for entry in data if entry[metric] >= 20 and entry[metric] <= 300)
        fail_count = sum(1 for entry in data if entry[metric] < 20 or entry[metric] > 300)
    elif metric == "First Pressure":
        pass_count = sum(1 for entry in data if entry[metric] >= 800 and entry[metric] <= 1100)
        fail_count = sum(1 for entry in data if entry[metric] < 800 or entry[metric] > 1100)
    elif metric == "First AccX":
        pass_count = sum(1 for entry in data if abs(entry[metric]) >= 0.0 and abs(entry[metric]) <= 0.2)
        fail_count = sum(1 for entry in data if abs(entry[metric]) < 0.0 or abs(entry[metric]) > 0.2)
    elif metric == "First AccY":
        pass_count = sum(1 for entry in data if abs(entry[metric]) >= 0.0 and abs(entry[metric]) <= 0.2)
        fail_count = sum(1 for entry in data if abs(entry[metric]) < 0.0 or abs(entry[metric]) > 0.2)
    elif metric == "First AccZ":
        pass_count = sum(1 for entry in data if entry[metric] and 0.8 <= abs(entry[metric]) <= 1.2)
        fail_count = sum(1 for entry in data if entry[metric] and not 0.8 <= abs(entry[metric]) <= 1.2)
    elif metric == "Temp Sensor":
        pass_count = sum(1 for entry in data if entry[metric] >= 20 and entry[metric] <= 30)
        fail_count = sum(1 for entry in data if entry[metric] < 20 or entry[metric] > 30)
    elif metric == "WiFi Scan":
        pass_count = sum(1 for entry in data if entry[metric] >= 1)
        fail_count = sum(1 for entry in data if entry[metric] < 1)
    else:
        pass_count = sum(1 for entry in data)
        fail_count = 0
    total_count = pass_count + fail_count
    pass_percentage = (pass_count / total_count) * 100 if total_count else 0
    fail_percentage = (fail_count / total_count) * 100 if total_count else 0
    return total_count, pass_percentage, fail_percentage, pass_count, fail_count

# Do statistical analysis for a given metric
def calculate_statistics(data, metric):
    values = [entry[metric] for entry in data if metric in entry and isinstance(entry[metric], (int, float))]
    if values:
        count = len(values)
        mean = np.mean(values)
        std = np.std(values)
        min_val = min(values)
        max_val = max(values)
    else:
        count = mean = std = min_val = max_val = 'N/A'  # In case there are no valid values
    return count, mean, std, min_val, max_val

# Function to generate and save the bar plot
def generate_bar_plot(data, metrics, image_path):
    for metric in metrics:
        if metric == "GPIB Current Reading":
            pass_count = sum(1 for entry in data if entry[metric] >= 430 and entry[metric] <= 860)
            fail_count = sum(1 for entry in data if entry[metric] < 430 or entry[metric] > 860)
        elif metric == "Battery voltage":
            pass_count = sum(1 for entry in data if entry[metric] >= 4600 and entry[metric] <= 6200)
            fail_count = sum(1 for entry in data if entry[metric] < 4600 or entry[metric] > 6200)
        elif metric == "Satellite":
            pass_count = sum(1 for entry in data if entry[metric] >= 1)
            fail_count = sum(1 for entry in data if entry[metric] < 1)
        elif metric == "Min Light":
            pass_count = sum(1 for entry in data if entry[metric] <= 5)
            fail_count = sum(1 for entry in data if entry[metric] > 5)
        elif metric == "Max Light":
            pass_count = sum(1 for entry in data if entry[metric] >= 20 and entry[metric] <= 300)
            fail_count = sum(1 for entry in data if entry[metric] < 20 or entry[metric] > 300)
        elif metric == "First Pressure":
            pass_count = sum(1 for entry in data if entry[metric] >= 800 and entry[metric] <= 1100)
            fail_count = sum(1 for entry in data if entry[metric] < 800 or entry[metric] > 1100)
        elif metric == "First AccX":
            pass_count = sum(1 for entry in data if abs(entry[metric]) >= 0.0 and abs(entry[metric]) <= 0.2)
            fail_count = sum(1 for entry in data if abs(entry[metric]) < 0.0 or abs(entry[metric]) > 0.2)
        elif metric == "First AccY":
            pass_count = sum(1 for entry in data if abs(entry[metric]) >= 0.0 and abs(entry[metric]) <= 0.2)
            fail_count = sum(1 for entry in data if abs(entry[metric]) < 0.0 or abs(entry[metric]) > 0.2)
        elif metric == "First AccZ":
            pass_count = sum(1 for entry in data if abs(entry[metric]) >= 0.8 and abs(entry[metric]) <= 1.2)
            fail_count = sum(1 for entry in data if abs(entry[metric]) > 1.2 or abs(entry[metric]) < 0.8)
        elif metric == "Temp Sensor":
            pass_count = sum(1 for entry in data if entry[metric] >= 20 and entry[metric] <= 30)
            fail_count = sum(1 for entry in data if entry[metric] < 20 or entry[metric] > 30)
        elif metric == "WiFi Scan":
            pass_count = sum(1 for entry in data if entry[metric] >= 1)
            fail_count = sum(1 for entry in data if entry[metric] < 1)
        else:
            pass_count = sum(1 for entry in data)
            fail_count = 0

    # Generate the bar plot
    x = np.arange(len(metrics))  # the label locations
    width = 0.35  # the width of the bars

    fig, ax = plt.subplots()
    rects1 = ax.bar(x - width/2, pass_count, width, label='Pass')
    rects2 = ax.bar(x + width/2, fail_count, width, label='Fail')

    # Add some text for labels, title and custom x-axis tick labels, etc.
    ax.set_ylabel('Counts')
    ax.set_title('Pass and Fail Counts by Metric')
    ax.set_xticks(x)
    ax.set_xticklabels(metrics, rotation=45, ha='right')
    ax.legend()

    # Save the plot as an image
    plt.tight_layout()
    plt.savefig(image_path)
    plt.close()

# Function to generate and save the histogram plot
def generate_histograms_for_metrics(data, metric, image_dir):
    # Extract values for the current metric, ensuring they are numeric and not missing
    values = [entry[metric] for entry in data]

    # Check if there are values to plot
    if values:
        # Generate the histogram plot
        plt.figure(figsize=(10, 6))  # Adjust the size as needed
        plt.hist(values, bins='auto', color='skyblue', alpha=0.7, rwidth=0.85)
        
        # Add labels and title
        plt.title(f'Histogram of {metric}')
        plt.xlabel(metric)
        plt.ylabel('Frequency')

        # Save the plot as an image
        image_path = os.path.join(image_dir, f'histogram_{metric}.png')
        plt.tight_layout()
        plt.savefig(image_path)
        plt.close()
        return image_path
    else:
        print(f"No valid data found for metric: {metric}")

# Divide the category
first_half = stats[0:5]
second_half = stats [5::]

def pass_fail_plot(headers, half_name):
    # Add a new slide for the stats table
    slide_layout = prs.slide_layouts[6]  # Assuming 5 is a layout that fits a table well
    summary_slide = prs.slides.add_slide(slide_layout)

    # Define table dimensions
    rows, cols = 6, len(headers) + 1  # Additional row for headers
    left, top, width, height = Inches(0.5), Inches(4.3), Inches(8.5), Inches(0.2)  # Modify as needed

    # Add a table to the slide (may need to adjust sizes and positions)
    table = summary_slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set rows headings
    table.cell(0, 0).text = 'Metric'
    table.cell(1, 0).text = 'Quantity'
    table.cell(2, 0).text = 'Pass %'
    table.cell(3, 0).text = 'Fail %'
    table.cell(4, 0).text = 'Pass'
    table.cell(5, 0).text = 'Fail'

    # Populate the table with data
    for i, metric in enumerate(headers, start=1):
        total, pass_percent, fail_percent, pass_count, fail_count = calculate_pass_fail(data, metric)
        table.cell(0, i).text = metric
        table.cell(1, i).text = str(total)
        table.cell(2, i).text = f"{pass_percent:.2f}%"
        table.cell(3, i).text = f"{fail_percent:.2f}%"
        table.cell(4, i).text = str(pass_count)
        table.cell(5, i).text = str(fail_count)

    # Define the path for the saved plot image
    plot_image_path = os.path.join(logdir, 'metrics_pass_fail_plot_'+half_name+'.png')

    # Generate and save the bar plot
    generate_bar_plot(data, headers, plot_image_path)

    # Insert the plot image into the slide
    left = Inches(1)
    top = Inches(0.1)
    summary_slide.shapes.add_picture(plot_image_path, left, top, width=Inches(8), height=Inches(4))

pass_fail_plot(first_half, "first_half")
pass_fail_plot(second_half, "second_half")

# Create a table and histogram plot for individual stats
for i in range(len(stats)):
    # Add a new slide for the summary table of statistics
    slide_layout = prs.slide_layouts[6]  # Choose a layout that fits a table well
    stats_slide = prs.slides.add_slide(slide_layout)

    # Define table dimensions
    rows, cols = 6, 2  # Additional row for headers
    left, top, width, height = Inches(3), Inches(4.3), Inches(4), Inches(0.2)  # Adjust as needed

    # Add a table to the slide
    table = stats_slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column headings
    table.cell(0, 0).text = 'Metric'
    table.cell(1, 0).text = 'Count'
    table.cell(2, 0).text = 'Mean'
    table.cell(3, 0).text = 'Std'
    table.cell(4, 0).text = 'Min'
    table.cell(5, 0).text = 'Max'

    # Populate the table with data
    count, mean, std, min_val, max_val = calculate_statistics(data, stats[i])
    table.cell(0, 1).text = stats[i]
    table.cell(1, 1).text = str(count)
    table.cell(2, 1).text = f"{mean:.2f}" if mean != 'N/A' else 'N/A'
    table.cell(3, 1).text = f"{std:.2f}" if std != 'N/A' else 'N/A'
    table.cell(4, 1).text = str(min_val)
    table.cell(5, 1).text = str(max_val)

    # Generate and save the bar plot
    histograms_dir = os.path.join(logdir, 'histograms')
    # Create the histograms folder if it doesn't exist
    if not os.path.exists(histograms_dir):
        os.makedirs(histograms_dir)
    image_path = generate_histograms_for_metrics(data, stats[i], histograms_dir)

    # Insert the plot image into the slide
    left = Inches(1)
    top = Inches(0.1)
    stats_slide.shapes.add_picture(image_path, left, top, width=Inches(8), height=Inches(4))

# Save the presentation
prs.save(logdir+"\\ASSY-MMI-Summary.pptx")

time.sleep(3)
subprocess.Popen(['start', '', logdir+"\\ASSY-MMI-Summary.pptx"], shell=True)
time.sleep(5)

print('ASSY-MMI done processing')