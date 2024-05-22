#!/usr/bin/env python
# coding: utf-8

import os
import re
import pandas as pd
from openpyxl import load_workbook
import openpyxl
from openpyxl.utils.dataframe import dataframe_to_rows
import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
import statsmodels.formula.api as stats

prs = Presentation()

def extract_reading(sensor_string, pattern):
    # Define the regular expression pattern to match the pressure reading
    match = re.search(pattern, sensor_string)

    if match:
        reading = match.group(1)
        return reading
    else:
        # Return None if no match is found
        return None

def time_value(string):
    return string.split("-")[0]+string.split("-")[1]+string.split("-")[2]+string.split("-")[3]+string.split("-")[4]+string.split("-")[5]

def date_stamp(string):
    return string.split("-")[0]+"/"+string.split("-")[1]+"/"+string.split("-")[2]+"-"

def clock_stamp(string):
    return string.split("-")[3]+":"+string.split("-")[4]+":"+string.split("-")[5]

def update_dict_in_array(array_of_dicts, new_object):
    for i, old_dict in enumerate(array_of_dicts):
        if old_dict["IMEI"] == new_object["IMEI"]:
            array_of_dicts[i] = new_object
            break

# Do statistical analysis for a given metric
def calculate_statistics(values):
    if values:
        count = len(values)
        mean = np.mean(values)
        std = np.std(values)
        min_val = min(values)
        max_val = max(values)
    else:
        count = mean = std = min_val = max_val = 'N/A'  # In case there are no valid values
    return count, mean, std, min_val, max_val

def FCT_dict(folder_name, lines, time_stamp, time_val):
    SN_MOB = None
    CCID = None
    IMEI = None
    #  NIST_Registers = None
    AccelX = None
    AccelY = None
    AccelZ = None 
    Pressure=  None 
    Temp = None 
    light = None
    wifi_scan = 0
    Button = False
    NTC = None
    Voltage = None
    Current = None
    wifi_version = None
    modem_version = None
    EEPROM1 = None
    EEPROM2 = None
    temp_offset = None
    EEPROM3 = None

    for line in lines:
        if re.search('CCID:\'', line):
            CCID = int(extract_reading(line, r'ID:\'(\d+)\''))
            SN_MOB = extract_reading(line, r'SN_MOB:\'(\w+)\'')
        if re.search('=IMEI:', line):
            IMEI = extract_reading(line, r"IMEI:(\d+)")
        if re.search('SSL3_', line):
            mcu_ver = line.split("\n")[0]
        if re.search('GSENSOR:', line):
            AccelX = float(re.findall(r'x\[(\-?\d+\.\d+)\]', line)[0])
            AccelY = float(re.findall(r'y\[(\-?\d+\.\d+)\]', line)[0])
            AccelZ = float(re.findall(r'z\[(\-?\d+\.\d+)\]', line)[0]) 
        if re.search('PRESS:', line):
            Pressure = float(re.findall(r'\d+\.\d+', line)[0])
        if re.search('TEMP:', line):
            Temp = float(extract_reading(line, r"\+TEMP:\[(\d+\.\d+)"))
        if re.search('LIGHT:', line):
            light = float(re.findall(r'\d+', line)[0])
        if re.search(r'Record (\d+): \+CWLAP',line):
            wifi_scan += 1
        if re.search('Button pushed',line):
            Button = True

        if folder_name == "Device_FCT":
            if re.search('get ntc adc value is', line):
                NTC = extract_reading(line, r"is(\d+\.\d+)")
            if re.search('	VBAT=', line):
                Voltage = float(re.search(r'\d+', line)[0])
            if re.search(r"\[DATARECV\]: \+.*", line):
                Current = float(re.search(r"\+?([-\d.]+)E", line).group(1))
        elif folder_name == "PCBA_FCT":
            if re.search('Voltage Regulator ', line):
                Voltage = float(re.findall(r'\d+', line)[0])
            if re.search('Bin version:', line):
                wifi_version = re.search(r'\d+\.\d+\.\d+\(ESP32C3-SPI\)', line).group()
            if re.search('BG', line):
                modem_version = line.split('\n')[0]
            if re.search('EEPROM2:', line):
                EEPROM1 = re.search(r'EEPROM1:\s*(0x[0-9A-Fa-f]+)', line).group(1)
                EEPROM2 = re.search(r'EEPROM2:\s*(0x[0-9A-Fa-f]+)', line).group(1)
                temp_offset = re.search(r'Temp Offset:\s*(0x[0-9A-Fa-f]+)', line).group(1)
                EEPROM3 = re.search(r'EEPROM3:\s*(0x[0-9A-Fa-f]+)', line).group(1)

    data_dict = {
        "IMEI": IMEI,
        "Timestamp": time_stamp,
        "CCID": str(CCID),
        "SN_MOB": SN_MOB,
        "MCU Version": mcu_ver,
        "Time Value": time_val,
        "AccelX": AccelX,
        "AccelY": AccelY,
        "AccelZ": AccelZ,
        "Pressure": Pressure,
        "Temp": Temp,
        "Light": light,
        "WiFi Scan Results": wifi_scan,
        "Button": Button,
        "Voltage (mV)": Voltage
    }
    if NTC:
        data_dict["NTC"] = NTC
    if Current:
        data_dict["Charge Current"] = Current
    if wifi_version:
        data_dict["Wifi Version"] = wifi_version
    if modem_version:
        data_dict["Modem version"] = modem_version
    if EEPROM1:
        data_dict["EEPROM1"] = EEPROM1
    if EEPROM2:
        data_dict["EEPROM2"] = EEPROM2
    if temp_offset:
        data_dict["Temp Offset"] = temp_offset
    if EEPROM3:
        data_dict["EEPROM3"] = EEPROM3


    return data_dict

def plot_FCT(folder_name):
    # Initialize data dictionaries
    data_dict = {}
    data = []
    logdir=os.getcwd()+"\\"+folder_name

    # Create Error Folder if not exists
    error_folder_path = os.path.join(logdir, "Error Folder")
    if not os.path.exists(error_folder_path):
        os.makedirs(error_folder_path)
    for file in os.listdir(logdir):
        if file != "Error Folder":
            time_stamp = date_stamp(file.split('_')[2])+clock_stamp(file.split('_')[2])
            time_val = int(time_value(file.split('_')[2]))
            file_path = os.path.join(logdir, file)
            f = open(file_path, 'r', encoding="utf8")
            try:
                lines = f.readlines()
                data_dict = FCT_dict(folder_name, lines, time_stamp, time_val)
                
                # Check for duplicate IDs
                contained = False
                replace = False
                for unit in data:
                    if unit["SN_MOB"] == data_dict["SN_MOB"]:
                        contained = True
                        if unit["Time Value"] < time_val:
                            replace = True
                if not contained:
                    data.append(data_dict)
                else:
                    if replace:
                        update_dict_in_array(data, data_dict)
            except Exception as e:
                f.close()
                # Move this file to directory "Error Folder"
                print(f"Error occurred on file: {file}. {str(e)}")
                error_file_path = os.path.join(error_folder_path, file)
                os.rename(file_path, error_file_path)
    print(len(data))

    # Create a dataframe
    df = pd.DataFrame(data)
    df = df[list(data[0].keys())]
    # Save dataframe as excel
    if os.path.isfile(os.getcwd()+'\\QDM065-Summary.xlsx') == False:
        df.to_excel(os.getcwd()+"\\QDM065-Summary.xlsx", index=False, sheet_name=folder_name)
    else:
        workbook = openpyxl.load_workbook(os.getcwd()+'\\QDM065-Summary.xlsx')  # load workbook if already exists
        sheet = workbook.create_sheet(folder_name)
        # append the dataframe results to the current excel file
        for row in dataframe_to_rows(df, header = True, index = False):
            sheet.append(row)
        workbook.save(os.getcwd()+'\\QDM065-Summary.xlsx')  # save workbook
        workbook.close()  # close workbook

    metrics_list = ["AccelX", "AccelY", "AccelZ", "Pressure", "Temp", "Light", "WiFi Scan Results", "Voltage (mV)"]
    if folder_name == "Device_FCT":
        metrics_list.append("Charge Current")
    for metric in metrics_list:
        values = [entry[metric] for entry in data]

        # Check if there are values to plot
        if values:
            # Add a new slide for the summary table of statistics
            slide_layout = prs.slide_layouts[6]  # Choose a layout that fits a table well
            stats_slide = prs.slides.add_slide(slide_layout)

            # Define table dimensions
            rows, cols = 6, 2  # Additional row for headers
            left, top, width, height = Inches(3), Inches(4.3), Inches(4), Inches(0.2)  # Adjust as needed

            # Add a table to the slide
            table = stats_slide.shapes.add_table(rows, cols, left, top, width, height).table

            # Set column headings
            table.cell(0, 0).text = str(metric)
            table.cell(1, 0).text = 'Count'
            table.cell(2, 0).text = 'Mean'
            table.cell(3, 0).text = 'Std'
            table.cell(4, 0).text = 'Min'
            table.cell(5, 0).text = 'Max'

            # Populate the table with data
            count, mean, std, min_val, max_val = calculate_statistics(values)
            table.cell(0, 1).text = str("Metric")
            table.cell(1, 1).text = str(count)
            table.cell(2, 1).text = f"{mean:.2f}" if mean != 'N/A' else 'N/A'
            table.cell(3, 1).text = f"{std:.2f}" if std != 'N/A' else 'N/A'
            table.cell(4, 1).text = str(min_val)
            table.cell(5, 1).text = str(max_val)

            # Generate the histogram plot
            plt.figure(figsize=(10, 6))  # Adjust the size as needed
            plt.hist(values, bins='auto', color='skyblue', alpha=0.7, rwidth=0.85)
            
            # Add labels and title
            plt.title(f'{str(metric)} for {folder_name}')
            plt.xlabel({str(metric)})
            plt.ylabel('Counts')

            # Save the plot as an image
            image_dir =  os.getcwd()+"\\images\\"
            if not os.path.exists(image_dir):
                os.makedirs(image_dir)
            image_path = os.path.join(image_dir, f'histogram_{folder_name}_{str(metric)}.png')
            plt.tight_layout()
            plt.savefig(image_path)
            plt.close()

            # Insert the plot image into the slide
            left = Inches(1)
            top = Inches(0.1)
            stats_slide.shapes.add_picture(image_path, left, top, width=Inches(8), height=Inches(4))
        else:
            print(f"No valid data found for metric: {folder_name}")

def plot_rf(folder_name):
    # Initialize data dictionaries
    data_dict = {}
    data = []
    logdir=os.getcwd()+"\\"+folder_name
    freq_list = []

    # Create Error Folder if not exists
    error_folder_path = os.path.join(logdir, "Error Folder")
    if not os.path.exists(error_folder_path):
        os.makedirs(error_folder_path)

    # Parsing data
    for file in os.listdir(logdir):
        if file != "Error Folder":
            time_stamp = date_stamp(file.split('_')[2])+clock_stamp(file.split('_')[2])
            time_val = int(time_value(file.split('_')[2]))
            file_path = os.path.join(logdir, file)
            f = open(file_path, 'r', encoding="utf8")
            IMEI = file.split('_')[0]
            temp_freq = -1
            temp_target = -1
            #freq_target_pair_list = []
            try:
                lines = f.readlines()
                for line in lines:
                    if re.search("HAN1;LTE;CONF:EARF:UL:cc1", line):
                        temp_freq = float(re.search(r"cc1\s+(\d+)", line).group(1)) / 10.0
                        if temp_freq not in freq_list:
                            freq_list.append(temp_freq)
                    if re.search("\'Test_LTE_TX_Power", line):
                        temp_target = float(extract_reading(line, r"'([\d.]+)'"))
                    if temp_freq != -1 and temp_target != -1:
                        #freq_target_pair_list.append((temp_freq, temp_target))
                        data_dict = {
                            "IMEI": IMEI,
                            "Timestamp": time_stamp,
                            "Time Value": time_val,
                            "Frequency": temp_freq,
                            "Measured Power": temp_target
                        }
                        data.append(data_dict)
                        temp_freq = -1
                        temp_target = -1
            except Exception as e:
                f.close()
                # Move this file to directory "Error Folder"
                print(f"Error occurred on file: {file}. {str(e)}")
                error_file_path = os.path.join(error_folder_path, file)
                os.rename(file_path, error_file_path)
    print(len(data))

    # Create a dataframe
    df = pd.DataFrame(data)
    df = df[list(data[0].keys())]
    # Save dataframe as excel
    if os.path.isfile(os.getcwd()+'\\QDM065-Summary.xlsx') == False:
        df.to_excel(os.getcwd()+"\\QDM065-Summary.xlsx", index=False, sheet_name=folder_name)
    else:
        workbook = openpyxl.load_workbook(os.getcwd()+'\\QDM065-Summary.xlsx')  # load workbook if already exists
        sheet = workbook.create_sheet(folder_name)
        # append the dataframe results to the current excel file
        for row in dataframe_to_rows(df, header = True, index = False):
            sheet.append(row)
        workbook.save(os.getcwd()+'\\QDM065-Summary.xlsx')  # save workbook
        workbook.close()  # close workbook

    for freq in freq_list:
        values = [entry["Measured Power"] for entry in data if entry["Frequency"] == freq]

        # Check if there are values to plot
        if values:
            # Add a new slide for the summary table of statistics
            slide_layout = prs.slide_layouts[6]  # Choose a layout that fits a table well
            stats_slide = prs.slides.add_slide(slide_layout)

            # Define table dimensions
            rows, cols = 6, 2  # Additional row for headers
            left, top, width, height = Inches(3), Inches(4.3), Inches(4), Inches(0.2)  # Adjust as needed

            # Add a table to the slide
            table = stats_slide.shapes.add_table(rows, cols, left, top, width, height).table

            # Set column headings
            table.cell(0, 0).text = 'Frequency'
            table.cell(1, 0).text = 'Count'
            table.cell(2, 0).text = 'Mean'
            table.cell(3, 0).text = 'Std'
            table.cell(4, 0).text = 'Min'
            table.cell(5, 0).text = 'Max'

            # Populate the table with data
            count, mean, std, min_val, max_val = calculate_statistics(values)
            table.cell(0, 1).text = str(freq)
            table.cell(1, 1).text = str(count)
            table.cell(2, 1).text = f"{mean:.2f}" if mean != 'N/A' else 'N/A'
            table.cell(3, 1).text = f"{std:.2f}" if std != 'N/A' else 'N/A'
            table.cell(4, 1).text = str(min_val)
            table.cell(5, 1).text = str(max_val)

            # Generate the histogram plot
            plt.figure(figsize=(10, 6))  # Adjust the size as needed
            plt.hist(values, bins='auto', color='skyblue', alpha=0.7, rwidth=0.85)
            
            # Add labels and title
            plt.title(f'{"Measured Power"} for {folder_name} at {str(freq)}')
            plt.xlabel("Measured Power")
            plt.ylabel('Counts')

            # Save the plot as an image
            image_dir =  os.getcwd()+"\\images\\"
            if not os.path.exists(image_dir):
                os.makedirs(image_dir)
            image_path = os.path.join(image_dir, f'histogram_{folder_name}_{str(freq)}.png')
            plt.tight_layout()
            plt.savefig(image_path)
            plt.close()

            # Insert the plot image into the slide
            left = Inches(1)
            top = Inches(0.1)
            stats_slide.shapes.add_picture(image_path, left, top, width=Inches(8), height=Inches(4))
        else:
            print(f"No valid data found for metric: {folder_name}")

def run_functions_safely():
    functions = [
        ("Device_FCT", plot_FCT),
        ("PCBA_FCT", plot_FCT),
        ("PCBA_FT_Conducted", plot_rf),
        ("PCBA_FT_Coupling", plot_rf)
    ]

    for arg, func in functions:
        try:
            # Add a slide with a title and content layout
            slide_layout = prs.slide_layouts[1]  # 0 is the layout for a title slide
            slide = prs.slides.add_slide(slide_layout)

            # Access the title and content placeholders
            title = slide.shapes.title
            title.text = arg
            title.text_frame.text = arg
            func(arg)
        except Exception as e:
            print(f"An error occurred while running {func.__name__}({arg}): {e}")
    
    prs.save(os.getcwd()+'\\charts.pptx')

run_functions_safely()
# plot_FCT("PCBA_FCT")
# plot_rf("PCBA_FT_Conducted")